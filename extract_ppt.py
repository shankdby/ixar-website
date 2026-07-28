import os
from pptx import Presentation

# Paths
ppt_path = r"c:\Users\sdube\OneDrive\Desktop\StudyFiles\superconductorinternshipfiles\ixar website\Website_Work_Copy.pptx"
output_dir = r"c:\Users\sdube\OneDrive\Desktop\StudyFiles\superconductorinternshipfiles\ixar website"
image_output_dir = os.path.join(output_dir, "extracted_images")

if not os.path.exists(image_output_dir):
    os.makedirs(image_output_dir)

prs = Presentation(ppt_path)
text_runs = []

print(f"Total slides found: {len(prs.slides)}")

for i, slide in enumerate(prs.slides):
    slide_num = i + 1
    slide_title = f"Slide {slide_num}"
    
    # Try to get slide title shape
    if slide.shapes.title:
        slide_title = slide.shapes.title.text.strip()
        
    text_runs.append(f"\n# {slide_title} (Slide {slide_num})\n")
    
    # Extract text from all text boxes
    slide_text = []
    image_count = 0
    
    for shape in slide.shapes:
        # Extract text
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                p_text = paragraph.text.strip()
                if p_text:
                    slide_text.append(p_text)
                    
        # Extract images
        if shape.shape_type == 13:  # 13 corresponds to MSO_SHAPE_TYPE.PICTURE
            try:
                image = shape.image
                image_bytes = image.blob
                ext = image.ext  # e.g., 'png', 'jpeg'
                
                image_count += 1
                image_filename = f"slide_{slide_num}_img_{image_count}.{ext}"
                image_path = os.path.join(image_output_dir, image_filename)
                
                with open(image_path, "wb") as f:
                    f.write(image_bytes)
            except Exception as e:
                pass
                
    if slide_text:
        text_runs.append("\n".join([f"- {text}" for text in slide_text]))
    else:
        text_runs.append("*(No text content)*")
        
    if image_count > 0:
        text_runs.append(f"\n*(Extracted {image_count} images from this slide)*")

# Save markdown output
md_path = os.path.join(output_dir, "extracted_content.md")
with open(md_path, "w", encoding="utf-8") as f:
    f.write("\n".join(text_runs))

print("Extraction completed successfully!")
print(f"Content saved to: {md_path}")
print(f"Images saved to: {image_output_dir}")
